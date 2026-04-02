"""PPT Builder Agent — Full LLM creative control over design.

Agent 3 works in two LLM steps:
  Step 1: DESIGN THEME — LLM creates a custom color palette, typography,
          layout style, and visual hierarchy tailored to the company/industry.
  Step 2: DESIGN SLIDES — LLM creates slides using the visual toolkit,
          with per-slide layout overrides, emphasis choices, and styling.
  Step 3: RENDER — Python executes the LLM's design decisions into PPTX.

The LLM has full control over:
  - Color palette (primary, secondary, accent, background, text)
  - Typography (font family, sizes for title/subtitle/body/caption)
  - Layout (header style, accent bar, footer, content margins)
  - Chart styling (series colors, show legends, label positions)
  - Per-slide overrides (different bg for title slide, emphasis colors)
"""

import json
import re
from io import BytesIO
from typing import Any

from google import genai
from google.genai import types
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from app.image_fetcher import fetch_company_logo
from app.settings import settings

# ── Slide dimensions ─────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ═══════════════════════════════════════════════════════════════
# Step 1: LLM designs the theme
# ═══════════════════════════════════════════════════════════════

THEME_PROMPT = """You are an elite presentation designer at a top consulting firm.
Design a CUSTOM visual theme for a PE due diligence presentation on: {company}

Consider the company's industry, brand, and personality when choosing colors.
A fintech company should feel different from a biotech or enterprise SaaS company.

Return STRICT JSON:
{{
  "theme_name": "descriptive name",
  "rationale": "Why these design choices fit this company",
  "colors": {{
    "bg_primary": "#hex",
    "bg_secondary": "#hex",
    "bg_card": "#hex",
    "accent_primary": "#hex",
    "accent_secondary": "#hex",
    "accent_danger": "#hex",
    "accent_success": "#hex",
    "text_primary": "#hex",
    "text_secondary": "#hex",
    "text_muted": "#hex",
    "header_bg": "#hex",
    "footer_bg": "#hex",
    "table_header_bg": "#hex",
    "table_row_odd": "#hex",
    "table_row_even": "#hex",
    "chart_colors": ["#hex1", "#hex2", "#hex3", "#hex4", "#hex5", "#hex6"]
  }},
  "typography": {{
    "font_family": "Calibri or Helvetica or Arial or Georgia",
    "title_size": 22,
    "subtitle_size": 13,
    "body_size": 15,
    "caption_size": 10,
    "metric_value_size": 26,
    "metric_label_size": 10
  }},
  "layout": {{
    "header_height": 0.55,
    "footer_height": 0.35,
    "accent_bar_width": 0.1,
    "content_left_margin": 0.65,
    "content_top_margin": 1.15,
    "gradient_angle": 135,
    "card_border_radius": "rounded",
    "slide_style": "modern_dark"
  }},
  "severity_colors": {{
    "high": "#hex",
    "medium_high": "#hex",
    "medium": "#hex",
    "low": "#hex"
  }}
}}

RULES:
- Background MUST be dark (readability on projectors)
- Text MUST be light/white for contrast
- Accent colors should feel related to the company's brand/industry
- Chart colors must be visually distinct from each other
- font_family must be a standard PPT font (Calibri, Arial, Helvetica, Georgia, Garamond)
- Sizes are in points (title: 20-28, body: 14-18, caption: 8-12)"""


def _design_theme(client, company: str) -> dict[str, Any]:
    """Ask LLM to design a custom theme."""
    prompt = THEME_PROMPT.format(company=company)
    response = client.models.generate_content(
        model=settings.gemini_slide_model,
        contents=prompt,
        config=types.GenerateContentConfig(
            temperature=0.4,  # Slightly creative
            response_mime_type="application/json",
        ),
    )
    try:
        return json.loads(response.text or "{}")
    except Exception:
        return {}


# ═══════════════════════════════════════════════════════════════
# Step 2: LLM designs the slides
# ═══════════════════════════════════════════════════════════════

SLIDE_PROMPT = """You are an elite PE presentation strategist at McKinsey.
Create a high-stakes IC due diligence deck for: {company}

DESIGN THEME you must follow:
{theme_json}

RESEARCH DATA:
{research_json}

You have a VISUAL TOOLKIT for each slide:
1. **chart**: {{"type": "pie|bar|column|doughnut", "categories": [...], "values": [numbers]}}
2. **table_data**: {{"headers": [...], "rows": [[...]]}}
3. **risk_blocks**: [{{"risk": "text", "severity": "High|Medium|Low"}}]

Return STRICT JSON:
{{
  "slides": [
    {{
      "slide_number": 1,
      "slide_type": "title|exec_summary|content|dashboard|two_column|sources",
      "title": "...",
      "subtitle": "...",
      "bullets": ["quantitative investor-grade bullets"],
      "key_stat": "most important insight for this slide",
      "source_ids": [1,2,3],
      "dashboard_metrics": [{{"label": "...", "value": "...", "trend": "up|down|flat"}}],
      "chart": null,
      "table_data": null,
      "risk_blocks": null,
      "style_overrides": {{
        "bg_override": "#hex or null",
        "accent_override": "#hex or null",
        "emphasis_color": "#hex or null",
        "title_size_override": null,
        "layout_hint": "full_width|split_left|split_right|centered"
      }}
    }}
  ]
}}

RULES:
- 14-20 slides
- Slide 1="title", slide 2="exec_summary", last="sources"
- 3+ dashboard slides with metrics AND charts
- USE CHARTS on 5+ slides: revenue=pie, KPIs=bar, funding=column, exit=doughnut, competition=column
- USE TABLES for: competitive positioning, comparable transactions
- USE risk_blocks for risk assessment
- style_overrides: use to vary visual emphasis across slides — not every slide should look identical
  - Title slide: maybe a different bg_override for impact
  - Dashboard slides: maybe a different accent color for data emphasis
  - Risk slide: emphasis_color could be the danger color
- layout_hint tells the renderer how to arrange content:
  - "full_width": bullets span full width
  - "split_left": chart/visual on right, bullets on left
  - "split_right": chart/visual on left, bullets on right
  - "centered": centered content (good for title, exec summary)
- Every non-title/sources slide: 3+ bullets, 2+ source_ids
- Make it VISUALLY IMMERSIVE — vary layouts, use the theme creatively"""


# ═══════════════════════════════════════════════════════════════
# Theme parser — converts LLM hex strings to RGBColor objects
# ═══════════════════════════════════════════════════════════════

def _hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert '#RRGGBB' to RGBColor. Fallback to white on error."""
    try:
        h = hex_str.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except Exception:
        return RGBColor(255, 255, 255)


class Theme:
    """Parsed theme from LLM output — every visual decision in one place."""

    def __init__(self, raw: dict[str, Any]):
        c = raw.get("colors", {})
        t = raw.get("typography", {})
        l = raw.get("layout", {})
        sev = raw.get("severity_colors", {})

        # Colors
        self.bg_primary = _hex_to_rgb(c.get("bg_primary", "#0A1628"))
        self.bg_secondary = _hex_to_rgb(c.get("bg_secondary", "#12233C"))
        self.bg_card = _hex_to_rgb(c.get("bg_card", "#142640"))
        self.accent_primary = _hex_to_rgb(c.get("accent_primary", "#C9A84C"))
        self.accent_secondary = _hex_to_rgb(c.get("accent_secondary", "#38BDA9"))
        self.accent_danger = _hex_to_rgb(c.get("accent_danger", "#DC5050"))
        self.accent_success = _hex_to_rgb(c.get("accent_success", "#38BDA9"))
        self.text_primary = _hex_to_rgb(c.get("text_primary", "#FFFFFF"))
        self.text_secondary = _hex_to_rgb(c.get("text_secondary", "#B4B4C3"))
        self.text_muted = _hex_to_rgb(c.get("text_muted", "#8896B0"))
        self.header_bg = _hex_to_rgb(c.get("header_bg", "#12233C"))
        self.footer_bg = _hex_to_rgb(c.get("footer_bg", "#12233C"))
        self.table_header_bg = _hex_to_rgb(c.get("table_header_bg", "#193255"))
        self.table_row_odd = _hex_to_rgb(c.get("table_row_odd", "#0E1C30"))
        self.table_row_even = _hex_to_rgb(c.get("table_row_even", "#12233C"))
        self.chart_colors = [_hex_to_rgb(h) for h in c.get("chart_colors", ["#C9A84C", "#38BDA9", "#DC5050", "#E6AA32", "#B4B4C3", "#8C64C8"])]

        # Typography
        self.font_family = t.get("font_family", "Calibri")
        self.title_size = Pt(min(max(t.get("title_size", 22), 18), 30))
        self.subtitle_size = Pt(min(max(t.get("subtitle_size", 13), 10), 18))
        self.body_size = Pt(min(max(t.get("body_size", 15), 12), 20))
        self.caption_size = Pt(min(max(t.get("caption_size", 10), 8), 13))
        self.metric_value_size = Pt(min(max(t.get("metric_value_size", 26), 20), 34))
        self.metric_label_size = Pt(min(max(t.get("metric_label_size", 10), 8), 13))

        # Layout
        self.header_h = Inches(min(max(l.get("header_height", 0.55), 0.4), 0.8))
        self.footer_h = Inches(min(max(l.get("footer_height", 0.35), 0.25), 0.5))
        self.accent_w = Inches(min(max(l.get("accent_bar_width", 0.1), 0.05), 0.2))
        self.content_left = Inches(min(max(l.get("content_left_margin", 0.65), 0.4), 1.0))
        self.content_top = Inches(min(max(l.get("content_top_margin", 1.15), 0.9), 1.5))
        self.gradient_angle = l.get("gradient_angle", 135)
        self.content_w = SLIDE_W - self.content_left - Inches(0.5)

        # Severity
        self.severity = {
            "high": _hex_to_rgb(sev.get("high", "#DC5050")),
            "medium-high": _hex_to_rgb(sev.get("medium_high", "#E68232")),
            "medium": _hex_to_rgb(sev.get("medium", "#E6AA32")),
            "medium-low": _hex_to_rgb(sev.get("medium_low", "#B4BE3C")),
            "low": _hex_to_rgb(sev.get("low", "#38BDA9")),
        }


def _default_theme() -> Theme:
    return Theme({})


# ═══════════════════════════════════════════════════════════════
# Chrome — applies theme to every slide
# ═══════════════════════════════════════════════════════════════

def _add_background(slide, theme: Theme, override_hex: str | None = None):
    bg = slide.background
    if override_hex:
        bg.fill.solid()
        bg.fill.fore_color.rgb = _hex_to_rgb(override_hex)
    else:
        bg.fill.gradient()
        bg.fill.gradient_angle = theme.gradient_angle
        stops = bg.fill.gradient_stops
        stops[0].color.rgb = theme.bg_primary
        stops[0].position = 0.0
        stops[1].color.rgb = theme.bg_secondary
        stops[1].position = 1.0


def _add_header(slide, theme: Theme, title: str, subtitle: str):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, theme.header_h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme.header_bg
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(theme.content_left, Inches(0.06), Inches(9.5), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = theme.title_size
    p.font.bold = True
    p.font.color.rgb = theme.text_primary
    p.font.name = theme.font_family
    if subtitle:
        sb = slide.shapes.add_textbox(theme.content_left, Inches(0.34), Inches(9.5), Inches(0.2))
        sp = sb.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = theme.subtitle_size
        sp.font.color.rgb = theme.accent_primary
        sp.font.italic = True
        sp.font.name = theme.font_family


def _add_accent(slide, theme: Theme):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, theme.header_h, theme.accent_w, SLIDE_H - theme.header_h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme.accent_primary
    bar.line.fill.background()


def _add_footer(slide, theme: Theme, company: str, num: int, total: int):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, SLIDE_H - theme.footer_h, SLIDE_W, theme.footer_h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme.footer_bg
    bar.line.fill.background()
    for text, x, align, color, bold in [
        (company, theme.content_left, PP_ALIGN.LEFT, theme.text_muted, False),
        (f"{num} / {total}", Inches(5.5), PP_ALIGN.CENTER, theme.text_muted, False),
        ("CONFIDENTIAL", Inches(10.0), PP_ALIGN.RIGHT, theme.accent_primary, True),
    ]:
        tb = slide.shapes.add_textbox(x, SLIDE_H - theme.footer_h + Inches(0.06), Inches(3.0), Inches(0.22))
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(9)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = theme.font_family
        p.alignment = align


def _add_chrome(slide, theme, company, title, subtitle, num, total, bg_override=None):
    _add_background(slide, theme, bg_override)
    _add_header(slide, theme, title, subtitle)
    _add_accent(slide, theme)
    _add_footer(slide, theme, company, num, total)


# ═══════════════════════════════════════════════════════════════
# Visual toolkit — chart/table/risk/progress helpers
# ═══════════════════════════════════════════════════════════════

def _style_chart(chart, theme: Theme):
    chart.has_legend = False
    if hasattr(chart, "category_axis"):
        ax = chart.category_axis
        ax.tick_labels.font.size = theme.caption_size
        ax.tick_labels.font.color.rgb = theme.text_secondary
        ax.tick_labels.font.name = theme.font_family
        ax.format.line.fill.background()
        ax.has_major_gridlines = False
    if hasattr(chart, "value_axis"):
        ax = chart.value_axis
        ax.tick_labels.font.size = theme.caption_size
        ax.tick_labels.font.color.rgb = theme.text_secondary
        ax.tick_labels.font.name = theme.font_family
        ax.format.line.fill.background()
        ax.has_major_gridlines = False


def _add_chart(slide, theme: Theme, spec: dict, x, y, w, h):
    """Render any chart from a spec dict."""
    ctype = spec.get("type", "bar")
    categories = spec.get("categories", [])
    values = spec.get("values", [])
    series = spec.get("series", [])
    if not categories:
        return

    chart_data = CategoryChartData()
    chart_data.categories = categories
    xl_type_map = {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "pie": XL_CHART_TYPE.PIE,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
    }
    xl_type = xl_type_map.get(ctype, XL_CHART_TYPE.BAR_CLUSTERED)

    if series:
        for s in series:
            chart_data.add_series(s.get("name", ""), s.get("values", []))
    else:
        chart_data.add_series("Value", values or [1] * len(categories))

    try:
        graphic = slide.shapes.add_chart(xl_type, x, y, w, h, chart_data)
        chart = graphic.chart
        _style_chart(chart, theme)
        # Color series/points using theme colors
        clrs = theme.chart_colors
        if ctype in ("pie", "doughnut"):
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.font.size = theme.caption_size
            chart.legend.font.color.rgb = theme.text_secondary
            chart.legend.font.name = theme.font_family
            chart.legend.include_in_layout = False
            plot = chart.plots[0]
            plot.has_data_labels = True
            plot.data_labels.font.size = theme.caption_size
            plot.data_labels.font.color.rgb = theme.text_primary
            plot.data_labels.font.bold = True
            for i, pt in enumerate(chart.series[0].points):
                pt.format.fill.solid()
                pt.format.fill.fore_color.rgb = clrs[i % len(clrs)]
        else:
            for i, s in enumerate(chart.series):
                s.format.fill.solid()
                s.format.fill.fore_color.rgb = clrs[i % len(clrs)]
            if len(chart.series) > 1:
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.font.size = theme.caption_size
                chart.legend.font.color.rgb = theme.text_secondary
                chart.legend.font.name = theme.font_family
    except Exception:
        pass


def _add_table(slide, theme: Theme, spec: dict, x, y, w):
    headers = spec.get("headers", [])
    rows = spec.get("rows", [])
    if not headers or not rows:
        return
    n_rows = len(rows) + 1
    n_cols = len(headers)
    total_h = Inches(0.38) * n_rows
    try:
        tbl_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, total_h)
        table = tbl_shape.table
        col_w = int(w / n_cols)
        for i in range(n_cols):
            table.columns[i].width = col_w
        for j, hdr in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = hdr
            cell.fill.solid()
            cell.fill.fore_color.rgb = theme.table_header_bg
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = theme.accent_primary
            p.font.name = theme.font_family
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for i, row in enumerate(rows[:8]):
            bg = theme.table_row_odd if i % 2 == 0 else theme.table_row_even
            for j, val in enumerate(row):
                if j >= n_cols:
                    break
                cell = table.cell(i + 1, j)
                cell.text = str(val or "")[:50]
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg
                p = cell.text_frame.paragraphs[0]
                p.font.size = theme.caption_size
                p.font.color.rgb = theme.text_primary
                p.font.name = theme.font_family
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        table.first_row = True
    except Exception:
        pass


def _add_risk_blocks(slide, theme: Theme, blocks: list[dict], x, y):
    bw = Inches(3.7)
    bh = Inches(1.0)
    gap = Inches(0.2)
    for i, block in enumerate(blocks[:6]):
        col = i % 3
        row = i // 3
        bx = x + (bw + gap) * col
        by = y + (bh + gap) * row
        sev = block.get("severity", "medium").lower()
        color = theme.severity.get(sev, theme.text_muted)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, bw, bh)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = block.get("risk", "")[:60]
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = theme.text_primary
        p.font.name = theme.font_family
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = sev.upper()
        p2.font.size = Pt(8)
        p2.font.color.rgb = theme.text_primary
        p2.font.name = theme.font_family
        p2.alignment = PP_ALIGN.CENTER


# ═══════════════════════════════════════════════════════════════
# Slide renderers — all theme-aware
# ═══════════════════════════════════════════════════════════════

def _bullets(slide, theme, items, x=None, y=None, w=None, max_n=6, size=None):
    x = x or theme.content_left
    y = y or theme.content_top
    w = w or theme.content_w
    fs = size or theme.body_size
    for b in items[:max_n]:
        tb = slide.shapes.add_textbox(x, y, w, Inches(0.4))
        tb.text_frame.word_wrap = True
        p = tb.text_frame.paragraphs[0]
        p.text = f"\u25B8  {b}"
        p.font.size = fs
        p.font.color.rgb = theme.text_primary
        p.font.name = theme.font_family
        y += Inches(0.42)
    return y


def _key_stat(slide, theme, text, x=None, y=None, w=None, emphasis=None):
    if not text:
        return
    x = x or theme.content_left
    y = y or Inches(5.3)
    w = w or Inches(11.5)
    color = _hex_to_rgb(emphasis) if emphasis else theme.accent_primary
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.5))
    box.fill.solid()
    box.fill.fore_color.rgb = theme.bg_card
    box.line.color.rgb = color
    box.line.width = Pt(1.2)
    p = box.text_frame.paragraphs[0]
    p.text = f"  \u2B50  {text}"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = theme.font_family


def _source_line(slide, theme, sids):
    if not sids:
        return
    tb = slide.shapes.add_textbox(theme.content_left, SLIDE_H - theme.footer_h - Inches(0.35), Inches(10), Inches(0.2))
    p = tb.text_frame.paragraphs[0]
    p.text = "Sources: " + "  ".join(f"[{s}]" for s in sids)
    p.font.size = Pt(9)
    p.font.color.rgb = theme.text_muted
    p.font.name = theme.font_family


def _render_title(slide, theme, company, data, total):
    overrides = data.get("style_overrides") or {}
    _add_background(slide, theme, overrides.get("bg_override"))
    _add_accent(slide, theme)
    _add_footer(slide, theme, company, 1, total)
    # Top accent line
    deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.08))
    deco.fill.solid()
    deco.fill.fore_color.rgb = theme.accent_primary
    deco.line.fill.background()
    # Company logo (if available)
    try:
        logo_path = fetch_company_logo(company)
        if logo_path:
            slide.shapes.add_picture(logo_path, Inches(5.65), Inches(1.1), width=Inches(2.0))
    except Exception:
        pass  # Skip logo on error

    # Company name
    tb = slide.shapes.add_textbox(Inches(1.5), Inches(2.3), Inches(10.5), Inches(1.2))
    p = tb.text_frame.paragraphs[0]
    p.text = company
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = theme.text_primary
    p.font.name = theme.font_family
    p.alignment = PP_ALIGN.CENTER
    # Subtitle
    sb = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(10.5), Inches(0.5))
    sp = sb.text_frame.paragraphs[0]
    sp.text = "Private Equity Due Diligence"
    sp.font.size = Pt(24)
    sp.font.color.rgb = theme.accent_primary
    sp.font.name = theme.font_family
    sp.alignment = PP_ALIGN.CENTER
    # Divider
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.1), Inches(4.5), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = theme.accent_primary
    line.line.fill.background()
    # Confidential
    nb = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(9.5), Inches(0.4))
    np_ = nb.text_frame.paragraphs[0]
    np_.text = "CONFIDENTIAL \u2014 For Investment Committee Use Only"
    np_.font.size = Pt(12)
    np_.font.color.rgb = theme.text_muted
    np_.font.name = theme.font_family
    np_.alignment = PP_ALIGN.CENTER
    ks = data.get("key_stat", "")
    if ks:
        kb = slide.shapes.add_textbox(Inches(2), Inches(5.2), Inches(9.5), Inches(0.35))
        kp = kb.text_frame.paragraphs[0]
        kp.text = ks
        kp.font.size = Pt(14)
        kp.font.color.rgb = theme.text_secondary
        kp.font.name = theme.font_family
        kp.alignment = PP_ALIGN.CENTER


def _render_content(slide, theme, company, data, num, total):
    overrides = data.get("style_overrides") or {}
    _add_chrome(slide, theme, company, data.get("title", ""), data.get("subtitle", ""), num, total, overrides.get("bg_override"))
    chart_spec = data.get("chart")
    table_spec = data.get("table_data")
    risk_spec = data.get("risk_blocks")
    bullets = data.get("bullets", [])
    hint = (overrides.get("layout_hint") or "full_width").lower()
    emphasis = overrides.get("emphasis_color")

    if chart_spec and hint in ("split_left", "split_right"):
        if hint == "split_left":
            _bullets(slide, theme, bullets, w=Inches(5.5), max_n=5)
            _add_chart(slide, theme, chart_spec, Inches(6.5), theme.content_top, Inches(6.0), Inches(4.0))
        else:
            _add_chart(slide, theme, chart_spec, theme.content_left, theme.content_top, Inches(5.5), Inches(4.0))
            _bullets(slide, theme, bullets, x=Inches(6.5), w=Inches(6.0), max_n=5)
    elif chart_spec:
        _bullets(slide, theme, bullets[:3], max_n=3)
        _add_chart(slide, theme, chart_spec, theme.content_left, Inches(3.2), Inches(11.5), Inches(3.2))
    elif table_spec:
        fy = _bullets(slide, theme, bullets[:2], max_n=2)
        _add_table(slide, theme, table_spec, theme.content_left, fy + Inches(0.1), theme.content_w)
    elif risk_spec:
        _bullets(slide, theme, bullets[:2], max_n=2)
        _add_risk_blocks(slide, theme, risk_spec, theme.content_left, Inches(2.8))
    else:
        _bullets(slide, theme, bullets)

    _key_stat(slide, theme, data.get("key_stat", ""), emphasis=emphasis)
    _source_line(slide, theme, data.get("source_ids", []))


def _render_dashboard(slide, theme, company, data, num, total):
    overrides = data.get("style_overrides") or {}
    _add_chrome(slide, theme, company, data.get("title", ""), data.get("subtitle", ""), num, total, overrides.get("bg_override"))
    metrics = data.get("dashboard_metrics", [])[:6]
    chart_spec = data.get("chart")

    if metrics:
        cols = 3 if len(metrics) > 4 else 2
        cw, ch = Inches(3.6), Inches(1.3)
        gx, gy = Inches(0.3), Inches(0.2)
        for i, m in enumerate(metrics):
            x = theme.content_left + (cw + gx) * (i % cols)
            y = theme.content_top + (ch + gy) * (i // cols)
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
            card.fill.solid()
            card.fill.fore_color.rgb = theme.bg_card
            card.line.color.rgb = _hex_to_rgb(overrides.get("accent_override") or "") if overrides.get("accent_override") else RGBColor(30, 50, 80)
            card.line.width = Pt(0.8)
            # Label
            lb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.12), cw - Inches(0.4), Inches(0.18))
            lp = lb.text_frame.paragraphs[0]
            lp.text = str(m.get("label", "")).upper()
            lp.font.size = theme.metric_label_size
            lp.font.color.rgb = theme.accent_primary
            lp.font.bold = True
            lp.font.name = theme.font_family
            # Value
            vb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.4), cw - Inches(0.4), Inches(0.4))
            vp = vb.text_frame.paragraphs[0]
            vp.text = str(m.get("value", "N/A"))
            vp.font.size = theme.metric_value_size
            vp.font.bold = True
            vp.font.color.rgb = theme.text_primary
            vp.font.name = theme.font_family
            # Trend
            trend = str(m.get("trend", "flat"))
            arrow = "\u2191" if trend == "up" else ("\u2193" if trend == "down" else "\u2192")
            tcolor = theme.accent_success if trend == "up" else (theme.accent_danger if trend == "down" else theme.text_muted)
            ttb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.95), cw - Inches(0.4), Inches(0.2))
            tp = ttb.text_frame.paragraphs[0]
            tp.text = f"{arrow}  {trend.upper()}"
            tp.font.size = Pt(10)
            tp.font.bold = True
            tp.font.color.rgb = tcolor
            tp.font.name = theme.font_family

    if chart_spec:
        cy = theme.content_top + Inches(3.0) if metrics else theme.content_top
        _add_chart(slide, theme, chart_spec, theme.content_left, cy, Inches(11.5), Inches(3.2))
    _source_line(slide, theme, data.get("source_ids", []))


def _render_exec_summary(slide, theme, company, data, num, total):
    overrides = data.get("style_overrides") or {}
    _add_chrome(slide, theme, company, data.get("title", ""), data.get("subtitle", ""), num, total, overrides.get("bg_override"))
    ks = data.get("key_stat", "Key insight")
    emphasis = overrides.get("emphasis_color")
    box_color = _hex_to_rgb(emphasis) if emphasis else theme.accent_primary
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, theme.content_left, theme.content_top, Inches(4.0), Inches(2.5))
    box.fill.solid()
    box.fill.fore_color.rgb = theme.bg_card
    box.line.color.rgb = box_color
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "\u2B50  KEY INSIGHT"
    p.font.size = theme.metric_label_size
    p.font.color.rgb = box_color
    p.font.bold = True
    p.font.name = theme.font_family
    vp = tf.add_paragraph()
    vp.text = ks
    vp.font.size = Pt(16)
    vp.font.color.rgb = theme.text_primary
    vp.font.bold = True
    vp.font.name = theme.font_family
    _bullets(slide, theme, data.get("bullets", []), x=Inches(5.2), w=Inches(7.5), max_n=6, size=Pt(14))
    if data.get("chart"):
        _add_chart(slide, theme, data["chart"], theme.content_left, Inches(4.0), Inches(11.5), Inches(2.5))
    _source_line(slide, theme, data.get("source_ids", []))


def _render_two_column(slide, theme, company, data, num, total):
    overrides = data.get("style_overrides") or {}
    _add_chrome(slide, theme, company, data.get("title", ""), data.get("subtitle", ""), num, total, overrides.get("bg_override"))
    if data.get("table_data"):
        _bullets(slide, theme, data.get("bullets", [])[:2], max_n=2, size=Pt(14))
        _add_table(slide, theme, data["table_data"], theme.content_left, Inches(2.2), theme.content_w)
    else:
        b = data.get("bullets", [])
        mid = (len(b) + 1) // 2
        _bullets(slide, theme, b[:mid], w=Inches(5.8), max_n=6)
        _bullets(slide, theme, b[mid:], x=Inches(7.0), w=Inches(5.8), max_n=6)
    _key_stat(slide, theme, data.get("key_stat", ""))
    _source_line(slide, theme, data.get("source_ids", []))


def _render_sources(slide, theme, company, data, num, total):
    _add_chrome(slide, theme, company, "Sources Appendix", "Traceable evidence index", num, total)
    _bullets(slide, theme, data.get("bullets", []), max_n=8, size=Pt(14))
    tb = slide.shapes.add_textbox(theme.content_left, Inches(5.5), theme.content_w, Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = "Cross-reference all source IDs with the source panel for audit-ready citation trails."
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = theme.text_muted
    p.font.name = theme.font_family


_RENDERERS = {
    "title": _render_title,
    "exec_summary": _render_exec_summary,
    "dashboard": _render_dashboard,
    "two_column": _render_two_column,
    "sources": _render_sources,
    "content": _render_content,
}


# ═══════════════════════════════════════════════════════════════
# Slide normalization + quality
# ═══════════════════════════════════════════════════════════════

def _normalize_slides(slides):
    out = []
    for i, s in enumerate(slides, 1):
        out.append({
            "slide_number": i,
            "slide_type": s.get("slide_type") or "content",
            "title": s.get("title") or f"Slide {i}",
            "subtitle": s.get("subtitle") or "",
            "bullets": [b for b in (s.get("bullets") or []) if b][:8],
            "key_stat": s.get("key_stat") or "",
            "source_ids": [sid for sid in (s.get("source_ids") or []) if sid][:8],
            "dashboard_metrics": (s.get("dashboard_metrics") or [])[:6],
            "chart": s.get("chart"),
            "table_data": s.get("table_data"),
            "risk_blocks": s.get("risk_blocks"),
            "style_overrides": s.get("style_overrides") or {},
        })
    return out[:24]


def _evaluate(slides):
    issues, score = [], 100
    if len(slides) < 14:
        issues.append(f"Need 14+ slides ({len(slides)})"); score -= 20
    dash = sum(1 for s in slides if s.get("dashboard_metrics"))
    if dash < 3:
        issues.append(f"Need 3+ dashboards ({dash})"); score -= 15
    vis = sum(1 for s in slides if s.get("chart") or s.get("table_data") or s.get("risk_blocks"))
    if vis < 3:
        issues.append(f"Need 3+ visual slides ({vis})"); score -= 10
    for i, s in enumerate(slides, 1):
        if s.get("slide_type") in ("title", "sources"):
            continue
        if len(s.get("bullets", [])) < 2:
            issues.append(f"Slide {i}: 2+ bullets"); score -= 4
        if len(s.get("source_ids", [])) < 1:
            issues.append(f"Slide {i}: source_ids"); score -= 3
    return score >= 80 and not issues, issues, score


# ═══════════════════════════════════════════════════════════════
# Mock theme + slides
# ═══════════════════════════════════════════════════════════════

def _mock_theme(company: str) -> dict[str, Any]:
    return {
        "theme_name": f"{company} Dark Executive",
        "colors": {
            "bg_primary": "#0A1628", "bg_secondary": "#12233C", "bg_card": "#142640",
            "accent_primary": "#C9A84C", "accent_secondary": "#38BDA9",
            "accent_danger": "#DC5050", "accent_success": "#38BDA9",
            "text_primary": "#FFFFFF", "text_secondary": "#B4B4C3", "text_muted": "#8896B0",
            "header_bg": "#12233C", "footer_bg": "#12233C",
            "table_header_bg": "#193255", "table_row_odd": "#0E1C30", "table_row_even": "#12233C",
            "chart_colors": ["#C9A84C", "#38BDA9", "#DC5050", "#E6AA32", "#B4B4C3", "#8C64C8"],
        },
        "typography": {"font_family": "Calibri", "title_size": 22, "subtitle_size": 13, "body_size": 15, "caption_size": 10, "metric_value_size": 26, "metric_label_size": 10},
        "layout": {"header_height": 0.55, "footer_height": 0.35, "accent_bar_width": 0.1, "content_left_margin": 0.65, "content_top_margin": 1.15, "gradient_angle": 135},
        "severity_colors": {"high": "#DC5050", "medium_high": "#E68232", "medium": "#E6AA32", "low": "#38BDA9"},
    }


def _mock_slides(company, research):
    """Generate mock slides — delegates to the existing mock logic."""
    from app.workspace import get_charts, get_tables, get_risks
    run_id = research.get("_run_id", "")

    def _g(d, k, default=""):
        return d.get(k, default) if isinstance(d, dict) else default

    def _gl(d, k):
        return d.get(k, []) if isinstance(d, dict) else []

    dm = research.get("dashboard_metrics", {}).get("metrics", [])
    profile = research.get("company_profile", {})
    mgmt = research.get("management_team", {})
    product = research.get("product_and_technology", {})
    bm = research.get("business_model", {})
    ue = research.get("unit_economics", {})
    fin = research.get("financial_signals", {})
    ml = research.get("market_landscape", {})
    cp = research.get("competitive_positioning", {})
    ce = research.get("customer_evidence", {})
    risk = research.get("risk_assessment", {})
    cat = research.get("catalysts_and_outlook", {})
    iv = research.get("investment_view", {})
    comps = research.get("comparable_transactions", {})
    exit_d = research.get("exit_analysis", {})
    sids = [s["id"] for s in research.get("all_sources", [])[:6]] or [1, 2, 3]

    risk_blocks = [{"risk": r.get("risk", "")[:60], "severity": r.get("severity", "Medium")} for r in _gl(risk, "key_risks")[:4] if isinstance(r, dict)]

    return {"slides": [
        {"slide_number": 1, "slide_type": "title", "title": company, "subtitle": "", "bullets": [], "key_stat": _g(fin, "arr_trajectory", ""), "source_ids": sids, "dashboard_metrics": [], "style_overrides": {"layout_hint": "centered"}},
        {"slide_number": 2, "slide_type": "exec_summary", "title": "Executive Summary", "subtitle": "Investment thesis", "bullets": [_g(profile, "summary"), _g(iv, "base_case"), _g(iv, "recommendation", "Proceed")], "key_stat": _g(iv, "summary", "Compelling opportunity"), "source_ids": sids, "dashboard_metrics": [], "chart": {"type": "bar", "categories": ["ARR", "Growth", "NRR", "Margin"], "values": [1000, 80, 150, 55]}},
        {"slide_number": 3, "slide_type": "content", "title": "Company Profile", "subtitle": f"Founded {_g(profile, 'founded')}", "bullets": [_g(profile, "summary"), f"Employees: {_g(profile, 'employee_estimate')}", *_gl(profile, "key_facts")[:3]], "key_stat": f"Employees: {_g(profile, 'employee_estimate')}", "source_ids": sids, "dashboard_metrics": [], "style_overrides": {"layout_hint": "full_width"}},
        {"slide_number": 4, "slide_type": "content", "title": "Management Team", "subtitle": "Leadership depth", "bullets": [_g(mgmt, "summary"), *[f"{e.get('name','')}: {e.get('background','')[:70]}" for e in _gl(mgmt, 'executives')[:3]]], "key_stat": _g(mgmt, "key_man_risk"), "source_ids": sids, "dashboard_metrics": []},
        {"slide_number": 5, "slide_type": "content", "title": "Product & Technology", "subtitle": "Moat analysis", "bullets": [_g(product, "summary"), *_gl(product, "core_offerings")[:3]], "key_stat": _g(product, "moat_hypothesis", "")[:80], "source_ids": sids, "dashboard_metrics": [], "style_overrides": {"layout_hint": "split_left"}, "chart": {"type": "bar", "categories": ["Speed", "Context", "Safety", "Compliance"], "values": [95, 90, 98, 92]}},
        {"slide_number": 6, "slide_type": "content", "title": "Business Model", "subtitle": "Revenue mechanics", "bullets": [_g(bm, "summary"), f"Pricing: {_g(bm, 'pricing_motion')}", *_gl(bm, "customer_segments")[:2]], "key_stat": "Land-and-expand with 140%+ NRR", "source_ids": sids, "dashboard_metrics": [], "chart": {"type": "pie", "categories": ["Subscriptions", "Usage-Based", "Services"], "values": [55, 40, 5]}, "style_overrides": {"layout_hint": "split_left"}},
        {"slide_number": 7, "slide_type": "dashboard", "title": "Unit Economics & KPIs", "subtitle": "Performance indicators", "bullets": [], "key_stat": "", "source_ids": sids, "dashboard_metrics": dm[:6], "chart": {"type": "bar", "categories": ["NRR", "Margin", "Rule of 40", "Retention", "LTV:CAC"], "values": [150, 55, 100, 94, 6]}},
        {"slide_number": 8, "slide_type": "dashboard", "title": "Financial Dashboard", "subtitle": "Traction signals", "bullets": [], "key_stat": _g(fin, "arr_trajectory"), "source_ids": sids, "dashboard_metrics": dm[:6], "chart": {"type": "column", "categories": ["Series A", "Series B", "Series C", "Series D"], "values": [124, 450, 2000, 4000]}},
        {"slide_number": 9, "slide_type": "content", "title": "Market Landscape", "subtitle": _g(ml, "market_size_estimate", "")[:50], "bullets": [_g(ml, "summary"), *_gl(ml, "industry_trends")[:3]], "key_stat": _g(ml, "market_size_estimate"), "source_ids": sids, "dashboard_metrics": []},
        {"slide_number": 10, "slide_type": "two_column", "title": "Competitive Positioning", "subtitle": "Differentiation", "bullets": [_g(cp, "summary"), *_gl(cp, "key_differentiators")[:2]], "key_stat": "Differentiated on safety + enterprise trust", "source_ids": sids, "dashboard_metrics": [], "table_data": {"headers": ["Competitor", "Strengths", "Weaknesses", "Our Advantage"], "rows": [[p.get("competitor", ""), p.get("strengths", "")[:35], p.get("weaknesses", "")[:35], p.get("vs_target", p.get("relative_position", ""))[:35]] for p in _gl(cp, "positioning_matrix")[:4] if isinstance(p, dict)]} if _gl(cp, "positioning_matrix") else None},
        {"slide_number": 11, "slide_type": "content", "title": "Comparable Transactions", "subtitle": "Valuation benchmarks", "bullets": [_g(comps, "summary")], "key_stat": _g(comps, "implied_valuation_range"), "source_ids": sids, "dashboard_metrics": [], "table_data": {"headers": ["Company", "Type", "Valuation", "Multiple"], "rows": [[t.get("target", ""), "Private", t.get("deal_value", ""), t.get("revenue_multiple", "")] for t in _gl(comps, "private_transactions")[:3] if isinstance(t, dict)] + [[t.get("company", ""), "Public", t.get("ev_revenue_multiple", ""), t.get("growth_rate", "")] for t in _gl(comps, "public_comps")[:3] if isinstance(t, dict)]} if _gl(comps, "private_transactions") or _gl(comps, "public_comps") else None},
        {"slide_number": 12, "slide_type": "content", "title": "Customer Evidence", "subtitle": "Traction signals", "bullets": [_g(ce, "summary"), f"Concentration: {_g(ce, 'concentration_risk')}", f"NPS: {_g(ce, 'nps_proxy')}", *_gl(ce, "case_studies")[:2]], "key_stat": _g(ce, "churn_signals"), "source_ids": sids, "dashboard_metrics": []},
        {"slide_number": 13, "slide_type": "dashboard", "title": "Risk Heatmap", "subtitle": "Downside analysis", "bullets": [_g(risk, "summary")], "key_stat": _g(risk, "overall_risk_rating"), "source_ids": sids, "dashboard_metrics": dm[:4], "risk_blocks": risk_blocks, "style_overrides": {"emphasis_color": "#DC5050"}},
        {"slide_number": 14, "slide_type": "content", "title": "Catalysts & Outlook", "subtitle": "Value creation levers", "bullets": [_g(cat, "summary"), *[f"{c.get('catalyst','')}: {c.get('impact','')}" for c in _gl(cat, "catalysts")[:3] if isinstance(c, dict)]], "key_stat": "Multiple catalysts within 12-24 months", "source_ids": sids, "dashboard_metrics": []},
        {"slide_number": 15, "slide_type": "content", "title": "Exit Analysis", "subtitle": "Returns profile", "bullets": [_g(exit_d, "summary"), *[f"{a.get('acquirer','')}: {a.get('rationale','')[:50]}" for a in _gl(exit_d, "strategic_acquirers")[:3] if isinstance(a, dict)]], "key_stat": f"Base IRR: {_g(_g(exit_d, 'implied_irr', {}), 'base_case')}", "source_ids": sids, "dashboard_metrics": [], "chart": {"type": "doughnut", "categories": ["Bear", "Base", "Bull"], "values": [25, 50, 25]}, "style_overrides": {"layout_hint": "split_left"}},
        {"slide_number": 16, "slide_type": "exec_summary", "title": "Investment Recommendation", "subtitle": "IC decision", "bullets": [_g(iv, "base_case"), _g(iv, "upside_case"), _g(iv, "downside_case")], "key_stat": _g(iv, "recommendation", "PROCEED"), "source_ids": sids, "dashboard_metrics": []},
        {"slide_number": 17, "slide_type": "sources", "title": "Sources Appendix", "subtitle": "Evidence index", "bullets": ["All evidence linked in source panel.", "Cross-reference before IC review.", "Flag data gaps as follow-ups."], "key_stat": "", "source_ids": sids, "dashboard_metrics": []},
    ]}


# ═══════════════════════════════════════════════════════════════
# Orchestration
# ═══════════════════════════════════════════════════════════════

def _slides_from_gemini(company, research, theme_raw):
    if settings.mock_mode or not settings.gemini_api_key:
        return _mock_slides(company, research)

    client = genai.Client(api_key=settings.gemini_api_key)
    prompt = SLIDE_PROMPT.format(
        company=company,
        theme_json=json.dumps(theme_raw, indent=2),
        research_json=json.dumps(research),
    )
    feedback = ""
    best, best_score = None, -1

    for _ in range(max(1, settings.slide_max_attempts)):
        p = prompt + (f"\nFix: {feedback}\n" if feedback else "")
        try:
            resp = client.models.generate_content(
                model=settings.gemini_slide_model, contents=p,
                config=types.GenerateContentConfig(temperature=0.2, response_mime_type="application/json"),
            )
            parsed = json.loads((resp.text or "{}").strip().strip("```json").strip("```").strip())
        except Exception:
            feedback = "Invalid JSON."
            continue
        if "slides" not in parsed:
            feedback = "Missing slides."
            continue
        parsed["slides"] = _normalize_slides(parsed["slides"])
        ok, issues, score = _evaluate(parsed["slides"])
        if score > best_score:
            best, best_score = parsed, score
        if ok:
            return parsed
        feedback = " ; ".join(issues)

    return best or _mock_slides(company, research)


def _enrich_from_workspace(slides, workspace):
    if not workspace:
        return slides
    charts = {c["chart_name"]: c for c in workspace.get("charts", [])}
    tables = {t["table_name"]: t for t in workspace.get("tables", [])}
    risks = workspace.get("risks", [])

    for s in slides:
        if s.get("chart") or s.get("table_data") or s.get("risk_blocks"):
            continue
        t = (s.get("title") or "").lower()
        # Charts
        for key, match in [("headline_metrics", "executive summary"), ("kpi_comparison", "unit economics"),
                           ("kpi_comparison", "kpi"), ("revenue_composition", "business model"),
                           ("funding_timeline", "financial"), ("exit_scenarios", "exit"),
                           ("competitor_comparison", "competitive")]:
            if match in t and key in charts:
                c = charts[key]
                s["chart"] = {"type": c["chart_type"], "categories": c["categories"], "values": c["values"]}
                break
        # Tables
        for key, match in [("competitive_positioning", "competitive"), ("comparable_transactions", "comparable"),
                           ("management_team", "management"), ("funding_rounds", "funding")]:
            if match in t and key in tables and not s.get("table_data"):
                tb = tables[key]
                s["table_data"] = {"headers": tb["headers"], "rows": tb["rows"]}
                break
        # Risks
        if "risk" in t and risks and not s.get("risk_blocks"):
            s["risk_blocks"] = [{"risk": r["risk"][:60], "severity": r["severity"]} for r in risks[:6]]

    return slides


def _create_pptx(company, slides, theme):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    total = len(slides)
    for data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        stype = data.get("slide_type", "content")
        renderer = _RENDERERS.get(stype, _render_content)
        if stype == "title":
            renderer(slide, theme, company, data, total)
        else:
            renderer(slide, theme, company, data, data["slide_number"], total)
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def build_presentation(company, research, workspace=None):
    # Step 1: Design theme
    if settings.mock_mode or not settings.gemini_api_key:
        theme_raw = _mock_theme(company)
    else:
        client = genai.Client(api_key=settings.gemini_api_key)
        theme_raw = _design_theme(client, company)
        if not theme_raw.get("colors"):
            theme_raw = _mock_theme(company)

    theme = Theme(theme_raw)

    # Step 2: Design slides
    payload = _slides_from_gemini(company, research, theme_raw)
    slides = _normalize_slides(payload["slides"])

    # Step 3: Enrich with workspace data
    slides = _enrich_from_workspace(slides, workspace)

    # Step 4: Render PPTX
    pptx = _create_pptx(company, slides, theme)
    return {"slides": slides}, pptx
